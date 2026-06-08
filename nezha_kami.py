#!/usr/bin/env python3
"""Nezha — Kami re-typeset: Nezha and His Moral Dilemma"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Kami Design System Constants ──

PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
OLIVE       = RGBColor(0x50, 0x4e, 0x49)
STONE       = RGBColor(0x6b, 0x6a, 0x64)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)

SERIF = "Source Han Serif SC"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Primitives ──

def blank_slide(prs, bg_color=PARCHMENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    return slide


def add_text(slide, text, left, top, width, height,
             font=SERIF, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_multiline(slide, lines, left, top, width, height,
                  font=SERIF, size=18, color=DARK_WARM,
                  align=PP_ALIGN.LEFT, line_spacing=Pt(28)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_comparison_items_col(slide, items, x_start, y_start, width, color):
    """Add bullet list items for comparison columns."""
    for i, item in enumerate(items):
        add_text(slide, f"•  {item}",
                 Inches(x_start), Inches(y_start + i * 0.6), Inches(width), Inches(0.55),
                 font=SERIF, size=16, color=color)


# ── Slide Templates ──

def cover_slide(prs, title, subtitle, footer=""):
    s = blank_slide(prs)
    add_line(s, Inches(1), Inches(2.0), Inches(11.33), weight_pt=0.5, color=OLIVE)
    add_text(s, title,
             Inches(1), Inches(2.3), Inches(11.33), Inches(2.0),
             font=SERIF, size=52, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.6), Inches(1), weight_pt=1.5)
    add_text(s, subtitle,
             Inches(1), Inches(4.9), Inches(11.33), Inches(0.6),
             font=SERIF, size=18, color=OLIVE, align=PP_ALIGN.CENTER)
    if footer:
        add_text(s, footer,
                 Inches(1), Inches(6.8), Inches(11.33), Inches(0.4),
                 font=SERIF, size=12, color=STONE, align=PP_ALIGN.CENTER)
    add_line(s, Inches(1), Inches(6.5), Inches(11.33), weight_pt=0.5, color=OLIVE)
    return s


def content_slide(prs, eyebrow, title, body_lines, page_num=None):
    s = blank_slide(prs)
    add_line(s, Inches(1.2), Inches(0.9), Inches(11), weight_pt=0.3, color=BORDER)
    add_text(s, eyebrow,
             Inches(1.2), Inches(1.1), Inches(10), Inches(0.4),
             font=SERIF, size=12, color=STONE)
    add_text(s, title,
             Inches(1.2), Inches(1.7), Inches(11.33), Inches(1.0),
             font=SERIF, size=34, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(2.8), Inches(3), weight_pt=1)
    add_multiline(s, body_lines,
                  Inches(1.2), Inches(3.3), Inches(11), Inches(3.5),
                  size=18, line_spacing=Pt(34))
    if page_num is not None:
        add_text(s, str(page_num).zfill(2),
                 Inches(12.3), Inches(6.9), Inches(1), Inches(0.3),
                 font=SERIF, size=10, color=STONE, align=PP_ALIGN.RIGHT)
    return s


def comparison_slide(prs, eyebrow, title,
                     left_title, left_items, right_title, right_items,
                     page_num=None):
    s = blank_slide(prs)
    add_line(s, Inches(1.2), Inches(0.9), Inches(11), weight_pt=0.3, color=BORDER)
    add_text(s, eyebrow,
             Inches(1.2), Inches(1.1), Inches(10), Inches(0.4),
             font=SERIF, size=12, color=STONE)
    add_text(s, title,
             Inches(1.2), Inches(1.7), Inches(11.33), Inches(0.8),
             font=SERIF, size=34, color=NEAR_BLACK)

    # Vertical divider
    divider = s.shapes.add_connector(1, Inches(6.67), Inches(2.5),
                                       Inches(6.67), Inches(6.5))
    divider.line.color.rgb = BORDER
    divider.line.width = Pt(1)

    # Left column (muted)
    add_line(s, Inches(1.2), Inches(2.7), Inches(5), weight_pt=0.5, color=BORDER)
    add_text(s, left_title,
             Inches(1.2), Inches(2.9), Inches(5), Inches(0.5),
             font=SERIF, size=20, color=OLIVE)
    add_comparison_items_col(s, left_items, 1.5, 3.7, 5, STONE)

    # Right column (full color)
    add_line(s, Inches(7.1), Inches(2.7), Inches(5), weight_pt=0.5, color=BORDER)
    add_text(s, right_title,
             Inches(7.1), Inches(2.9), Inches(5), Inches(0.5),
             font=SERIF, size=20, color=NEAR_BLACK)
    add_comparison_items_col(s, right_items, 7.4, 3.7, 5, DARK_WARM)

    if page_num is not None:
        add_text(s, str(page_num).zfill(2),
                 Inches(12.3), Inches(6.9), Inches(1), Inches(0.3),
                 font=SERIF, size=10, color=STONE, align=PP_ALIGN.RIGHT)
    return s


def quote_slide(prs, eyebrow, quote, source, page_num=None):
    s = blank_slide(prs)
    add_line(s, Inches(1.2), Inches(0.9), Inches(11), weight_pt=0.3, color=BORDER)
    add_text(s, eyebrow,
             Inches(1.2), Inches(1.1), Inches(10), Inches(0.4),
             font=SERIF, size=12, color=STONE)
    add_text(s, "\u201c",
             Inches(1.2), Inches(2.0), Inches(11.33), Inches(1.5),
             font=SERIF, size=96, color=BRAND, align=PP_ALIGN.LEFT)
    add_text(s, quote,
             Inches(1.2), Inches(2.8), Inches(10.93), Inches(2.8),
             font=SERIF, size=28, color=NEAR_BLACK,
             align=PP_ALIGN.LEFT, vanchor=MSO_ANCHOR.MIDDLE, italic=True)
    add_text(s, f"\u2014 {source}",
             Inches(1.2), Inches(5.8), Inches(10.93), Inches(0.4),
             font=SERIF, size=14, color=OLIVE, align=PP_ALIGN.RIGHT)
    if page_num is not None:
        add_text(s, str(page_num).zfill(2),
                 Inches(12.3), Inches(6.9), Inches(1), Inches(0.3),
                 font=SERIF, size=10, color=STONE, align=PP_ALIGN.RIGHT)
    return s


def ending_slide(prs, message, subtitle=""):
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, message,
             Inches(1.5), Inches(2.2), Inches(10.33), Inches(2.0),
             font=SERIF, size=36, color=WHITE, align=PP_ALIGN.CENTER)
    add_line(s, Inches(5.67), Inches(4.3), Inches(2), weight_pt=1.5, color=WHITE)
    if subtitle:
        add_text(s, subtitle,
                 Inches(1.5), Inches(4.6), Inches(10.33), Inches(0.6),
                 font=SERIF, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    return s


# ── Main ──

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="Nezha_Moral_Dilemma_kami.pptx")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Cover ──
    cover_slide(prs,
        title="Nezha and His\nMoral Dilemma",
        subtitle="哪吒 — A Story from Chinese Mythology",
        footer="A Pre-Class Presentation")

    # ── Slide 2: Who Is Nezha? ──
    content_slide(prs,
        eyebrow="01  ·  BACKGROUND",
        title="Who Is Nezha?",
        body_lines=[
            "Born from the Mo Wan (Lotus Root)",
            "People believed the Mo Wan would bring disaster",
            "Nezha was expected to become evil in the future",
            "Many people were afraid of him and did not trust him",
        ],
        page_num=2)

    # ── Slide 3: The Moral Dilemma ──
    comparison_slide(prs,
        eyebrow="02  ·  THE DILEMMA",
        title="The Moral Dilemma",
        left_title="Accept Fate",
        left_items=[
            "Accept what others thought about him",
            "Become the monster they expected",
            "Express his anger freely",
        ],
        right_title="Fight Fate",
        right_items=[
            "Fight against what others believed",
            "Prove that he was a good person",
            "Protect the innocent",
        ],
        page_num=3)

    # ── Slide 4: Two Paths, Both Difficult ──
    comparison_slide(prs,
        eyebrow="03  ·  TWO PATHS",
        title="Two Paths, Both Difficult",
        left_title="If He Chose Kindness",
        left_items=[
            "He would still face misunderstanding",
            "He would experience loneliness",
            "No one would trust him despite his good deeds",
        ],
        right_title="If He Chose Revenge",
        right_items=[
            "He could finally express his anger",
            "He would prove them right about him",
            "But innocent people would be hurt",
        ],
        page_num=4)

    # ── Slide 5: Nezha's Choice (Quote) ──
    quote_slide(prs,
        eyebrow="04  ·  THE CHOICE",
        quote="Nezha chose to protect others\ninstead of hurting them.",
        source="People are not defined by their birth,\nnor by what others think of them.",
        page_num=5)

    # ── Slide 6: Ending ──
    ending_slide(prs,
        message="Even when the world\nmisunderstands us,",
        subtitle="we should still choose what is right.\n\nThank You")

    prs.save(args.out)
    print(f"✅ Saved: {args.out}")


if __name__ == '__main__':
    main()
